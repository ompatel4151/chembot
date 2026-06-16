"""
ingest_free.py — Rebuild vectorstore from data/ folder.
Supports PDF, PPTX, PPT (via python-pptx), and DOCX.

Usage:
    python ingest_free.py

Output:
    vectorstore/index.faiss  (384-dim)
    vectorstore/docs.pkl
"""

import os
import pickle
import numpy as np

DATA_PATH  = "data"
VECTOR_DIR = "vectorstore"
INDEX_FILE = os.path.join(VECTOR_DIR, "index.faiss")
DOCS_FILE  = os.path.join(VECTOR_DIR, "docs.pkl")

CHUNK_SIZE = 900
OVERLAP    = 200


def chunk_text(text: str):
    text = " ".join(text.split())   # normalize whitespace
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start:start + CHUNK_SIZE].strip()
        if chunk:
            chunks.append(chunk)
        start += CHUNK_SIZE - OVERLAP
    return chunks


def extract_pdf(path: str):
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        txt = (page.extract_text() or "").strip()
        if txt:
            pages.append((i + 1, txt))
    return pages   # list of (page_num, text)


def extract_pptx(path: str):
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        combined = " ".join(texts).strip()
        if combined:
            pages.append((i + 1, combined))
    return pages


def extract_docx(path: str):
    from docx import Document
    doc = Document(path)
    text = " ".join(p.text for p in doc.paragraphs if p.text.strip())
    return [(1, text)] if text.strip() else []


def extract_file(path: str):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    if ext == ".docx":
        return extract_docx(path)
    return []


def main():
    os.makedirs(VECTOR_DIR, exist_ok=True)
    os.makedirs(DATA_PATH, exist_ok=True)

    supported = (".pdf", ".pptx", ".ppt", ".docx")
    files = [f for f in os.listdir(DATA_PATH) if f.lower().endswith(supported)]

    if not files:
        print(f"❌ No supported files found in '{DATA_PATH}/'.")
        print(f"   Add PDF, PPTX, or DOCX files and run again.")
        return

    docs = []
    for filename in sorted(files):
        path = os.path.join(DATA_PATH, filename)
        print(f"📄 {filename}")
        try:
            pages = extract_file(path)
        except Exception as e:
            print(f"   ⚠️  Skipped ({e})")
            continue

        for page_num, text in pages:
            for chunk_i, chunk in enumerate(chunk_text(text)):
                docs.append({
                    "text": chunk,
                    "meta": {
                        "source_file": filename,
                        "page_number": page_num,
                        "chunk_id":    f"{filename}_p{page_num}_c{chunk_i+1}"
                    }
                })

    print(f"\n✅ {len(docs)} chunks from {len(files)} file(s)")

    print("⏳ Loading embedding model (first run downloads ~90 MB)…")
    from sentence_transformers import SentenceTransformer
    model = SentenceTransformer("all-MiniLM-L6-v2")

    print("⏳ Embedding chunks…")
    BATCH = 64
    all_vecs = []
    for i in range(0, len(docs), BATCH):
        batch = [d["text"] for d in docs[i:i + BATCH]]
        vecs  = model.encode(batch, show_progress_bar=False)
        all_vecs.extend(vecs)
        print(f"   {min(i + BATCH, len(docs))}/{len(docs)}")

    emb = np.array(all_vecs, dtype=np.float32)

    import faiss
    dim   = emb.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(emb)
    faiss.write_index(index, INDEX_FILE)

    with open(DOCS_FILE, "wb") as f:
        pickle.dump(docs, f)

    print(f"\n✅ Vectorstore saved!")
    print(f"   {INDEX_FILE}  ({dim}-dim, {len(docs)} vectors)")
    print(f"   {DOCS_FILE}")
    print(f"\nRestart app.py and you're good to go 🚀")


if __name__ == "__main__":
    main()
